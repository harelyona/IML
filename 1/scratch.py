import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def generate_pv_presentation():
    """
        Generates a comprehensive PowerPoint presentation on Photovoltaic Technology
        based on the textbook chapter.
        """
    prs = Presentation()

    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Define common styles
    title_font_size = Pt(44)
    subtitle_font_size = Pt(32)
    body_font_size = Pt(24)
    caption_font_size = Pt(18)

    title_color = RGBColor(31, 73, 125)  # Dark blue
    subtitle_color = RGBColor(68, 114, 196)  # Medium blue
    accent_color = RGBColor(112, 173, 71)  # Green

    # --- Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Photovoltaic Power Technology"
    subtitle.text = "Principles, Applications, and Future Trends"

    title.text_frame.paragraphs[0].font.size = title_font_size
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    subtitle.text_frame.paragraphs[0].font.size = subtitle_font_size

    # --- Introduction Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "What is Photovoltaic Technology?"
    title.text_frame.paragraphs[0].font.color.rgb = title_color

    tf = content.text_frame
    p = tf.paragraphs[0]
    p.text = "• One of only two methods to generate significant electric power"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Discovered by Becquerel in 1839, developed as power source in 1954"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Converts sunlight directly into electricity with no moving parts"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Based on semiconductor physics and the photovoltaic effect"
    p.font.size = body_font_size

    # --- How PV Cells Work Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "How Photovoltaic Cells Work"
    title.text_frame.paragraphs[0].font.color.rgb = title_color

    tf = content.text_frame
    p = tf.paragraphs[0]
    p.text = "1. Photons from sunlight strike semiconductor material"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "2. Photons with sufficient energy create electron-hole pairs"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "3. Built-in electric field at p-n junction separates charges"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "4. Electrons flow through external circuit as electric current"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "5. Typical silicon cells produce ~0.5-0.6V"
    p.font.size = body_font_size

    # --- P-N Junction Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "The P-N Junction: Heart of a Solar Cell"
    title.text_frame.paragraphs[0].font.color.rgb = title_color

    tf = content.text_frame
    p = tf.paragraphs[0]
    p.text = "• P-type: Silicon doped with boron (has 'holes' or positive carriers)"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• N-type: Silicon doped with phosphorus (has excess electrons)"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Junction creates built-in electric field and depletion region"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Light creates electron-hole pairs that are separated by this field"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• This charge separation generates voltage and current"
    p.font.size = body_font_size

    # --- I-V Characteristics Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Electrical Characteristics"
    title.text_frame.paragraphs[0].font.color.rgb = title_color

    tf = content.text_frame
    p = tf.paragraphs[0]
    p.text = "• Short-circuit current (Isc): Maximum current when V=0"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Open-circuit voltage (Voc): Maximum voltage when I=0"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Fill factor: Ratio of maximum power to product of Isc and Voc"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Maximum power point: Optimal operating point on I-V curve"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Temperature effects: Power decreases as temperature increases"
    p.font.size = body_font_size

    # --- Efficiency Limitations Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Efficiency Limitations"
    title.text_frame.paragraphs[0].font.color.rgb = title_color

    tf = content.text_frame
    p = tf.paragraphs[0]
    p.text = "• Photons with energy less than band gap (~23% loss)"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Excess photon energy lost as heat (~33% loss)"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Voltage factor losses (~20% loss)"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Fill factor and collection inefficiencies (~22% combined)"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Theoretical limit (Shockley-Queisser): ~31% for single-junction"
    p.font.size = body_font_size

    # --- Cell Technologies Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "PV Cell Technologies"
    title.text_frame.paragraphs[0].font.color.rgb = title_color

    tf = content.text_frame
    p = tf.paragraphs[0]
    p.text = "First Generation: Crystalline Silicon (80% of market)"
    p.font.size = body_font_size
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "• Single-crystal Si (15-25% efficiency)"
    p.font.size = body_font_size
    p.level = 1

    p = tf.add_paragraph()
    p.text = "• Multi-crystalline Si (13-20% efficiency)"
    p.font.size = body_font_size
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Second Generation: Thin Film Technologies"
    p.font.size = body_font_size
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "• Amorphous Si, CdTe, CIGS (7-20% efficiency)"
    p.font.size = body_font_size
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Third Generation: Multi-junction cells, concentrators"
    p.font.size = body_font_size
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "• Up to 44% efficiency with concentration"
    p.font.size = body_font_size
    p.level = 1

    # --- System Applications Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "PV System Applications"
    title.text_frame.paragraphs[0].font.color.rgb = title_color

    tf = content.text_frame
    p = tf.paragraphs[0]
    p.text = "Grid-Connected Systems (97% of global installations)"
    p.font.size = body_font_size
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "• Rooftop installations (residential/commercial)"
    p.font.size = body_font_size
    p.level = 1

    p = tf.add_paragraph()
    p.text = "• Building-integrated PV (BIPV)"
    p.font.size = body_font_size
    p.level = 1

    p = tf.add_paragraph()
    p.text = "• Utility-scale solar farms"
    p.font.size = body_font_size
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Stand-Alone Systems"
    p.font.size = body_font_size
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "• Rural electrification in developing countries"
    p.font.size = body_font_size
    p.level = 1

    p = tf.add_paragraph()
    p.text = "• Telecommunications, water pumping, lighting"
    p.font.size = body_font_size
    p.level = 1

    # --- Market Growth Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Market Growth & Economics"
    title.text_frame.paragraphs[0].font.color.rgb = title_color

    tf = content.text_frame
    p = tf.paragraphs[0]
    p.text = "• Explosive growth: 200 MW in 1990 → 80+ GW in 2012"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Cost reduction: $40/W in 1980s → ~$1/W in 2013"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Grid parity achieved in many regions"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Positive feedback loop: Policy → Demand → Production → Lower costs"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Energy payback time: 2-3 years (temperate climate)"
    p.font.size = body_font_size

    # --- Environmental Benefits Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Environmental & Social Benefits"
    title.text_frame.paragraphs[0].font.color.rgb = title_color

    tf = content.text_frame
    p = tf.paragraphs[0]
    p.text = "• Zero emissions during operation"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• No noise pollution"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Long lifetime (20+ years guaranteed, much longer expected)"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Rural electrification in developing nations"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Energy independence and security"
    p.font.size = body_font_size

    # --- Future Outlook Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Future Outlook"
    title.text_frame.paragraphs[0].font.color.rgb = title_color

    tf = content.text_frame
    p = tf.paragraphs[0]
    p.text = "• Continued efficiency improvements"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Further cost reductions through manufacturing scale"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Building integration becoming standard"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• New materials and concepts (perovskites, organic PV)"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• PV becoming mainstream electricity generation worldwide"
    p.font.size = body_font_size

    # --- Conclusion Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Conclusion: PV Power Advantage"
    title.text_frame.paragraphs[0].font.color.rgb = title_color

    tf = content.text_frame
    p = tf.paragraphs[0]
    p.text = "• Universal applicability (solar radiation available everywhere)"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Modular nature (scalable from watts to megawatts)"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Reliability, long life, ease of use"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• No emissions, no noise, minimal maintenance"
    p.font.size = body_font_size

    p = tf.add_paragraph()
    p.text = "• Continuing improvements in technology and cost"
    p.font.size = body_font_size

    # --- Questions Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title

    title.text = "Questions?"
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Save the presentation
    output_file = "Photovoltaic_Technology_Presentation.pptx"
    prs.save(output_file)
    return output_file


if __name__ == "__main__":
    output_file = generate_pv_presentation()
    print(f"Successfully created presentation: {output_file}")

    # Display the absolute path to the file
    abs_path = os.path.abspath(output_file)
    print(f"Full path: {abs_path}")